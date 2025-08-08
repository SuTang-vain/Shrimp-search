"""
增强的文档管理器 v2.0
支持文档预处理缓存、增量更新和多格式处理
"""

import os
import json
import hashlib
import pickle
import time
from typing import List, Dict, Any, Optional, Set, Tuple
from pathlib import Path
import logging
from dataclasses import dataclass, asdict
from datetime import datetime

# 多模态处理器
try:
    from enhanced_multimodal_processor import EnhancedMultimodalProcessor
    MULTIMODAL_AVAILABLE = True
except ImportError:
    MULTIMODAL_AVAILABLE = False

# 文档处理库
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

@dataclass
class DocumentMetadata:
    """文档元数据"""
    file_path: str
    file_hash: str
    file_size: int
    last_modified: float
    processing_time: float
    chunk_count: int
    content_type: str
    format_type: str
    created_at: str
    version: str = "1.0"

@dataclass
class DocumentChunk:
    """文档块"""
    chunk_id: str
    content: str
    metadata: Dict[str, Any]
    embedding: Optional[List[float]] = None
    chunk_type: str = "text"  # text, image, table, code

class EnhancedDocumentManager:
    """增强的文档管理器"""
    
    def __init__(self, cache_dir: str = "document_cache", 
                 enable_cache: bool = True,
                 max_cache_size_mb: int = 1000):
        self.cache_dir = Path(cache_dir)
        self.enable_cache = enable_cache
        self.max_cache_size_mb = max_cache_size_mb
        
        # 创建缓存目录
        if self.enable_cache:
            self.cache_dir.mkdir(exist_ok=True)
            (self.cache_dir / "metadata").mkdir(exist_ok=True)
            (self.cache_dir / "chunks").mkdir(exist_ok=True)
            (self.cache_dir / "embeddings").mkdir(exist_ok=True)
        
        # 支持的文档格式
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.txt': self._process_txt,
            '.md': self._process_markdown,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_pptx,
            '.ppt': self._process_ppt,
            '.json': self._process_json,
            '.xml': self._process_xml,
            '.html': self._process_html,
            '.rtf': self._process_rtf
        }
        
        # 文档缓存索引
        self.document_index = self._load_document_index()
        
        # 初始化多模态处理器
        self.multimodal_processor = None
        if MULTIMODAL_AVAILABLE:
            try:
                self.multimodal_processor = EnhancedMultimodalProcessor()
                print("✅ 多模态处理器已启用")
            except Exception as e:
                print(f"⚠️ 多模态处理器初始化失败: {e}")
        
        # 设置日志
        self.logger = logging.getLogger(__name__)
        
        print("✅ 增强文档管理器初始化完成")
        print(f"📁 缓存目录: {self.cache_dir}")
        print(f"📋 支持格式: {list(self.supported_formats.keys())}")
        print(f"🔍 多模态处理: {'✓' if self.multimodal_processor else '✗'}")
    
    def _load_document_index(self) -> Dict[str, DocumentMetadata]:
        """加载文档索引"""
        index_file = self.cache_dir / "document_index.json"
        if index_file.exists():
            try:
                with open(index_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return {k: DocumentMetadata(**v) for k, v in data.items()}
            except Exception as e:
                self.logger.warning(f"加载文档索引失败: {e}")
        return {}
    
    def _save_document_index(self):
        """保存文档索引"""
        if not self.enable_cache:
            return
        
        index_file = self.cache_dir / "document_index.json"
        try:
            data = {k: asdict(v) for k, v in self.document_index.items()}
            with open(index_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            self.logger.error(f"保存文档索引失败: {e}")
    
    def _get_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            self.logger.error(f"计算文件哈希失败: {e}")
            return ""
    
    def _is_document_cached(self, file_path: str) -> bool:
        """检查文档是否已缓存且有效"""
        if not self.enable_cache or file_path not in self.document_index:
            return False
        
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                return False
            
            # 获取文件信息
            file_stat = os.stat(file_path)
            cached_metadata = self.document_index[file_path]
            
            # 检查文件是否被修改
            if (file_stat.st_mtime != cached_metadata.last_modified or
                file_stat.st_size != cached_metadata.file_size):
                return False
            
            # 检查缓存文件是否存在
            cache_file = self.cache_dir / "chunks" / f"{cached_metadata.file_hash}.pkl"
            return cache_file.exists()
            
        except Exception as e:
            self.logger.error(f"检查缓存状态失败: {e}")
            return False
    
    def _load_cached_chunks(self, file_path: str) -> List[DocumentChunk]:
        """加载缓存的文档块"""
        try:
            metadata = self.document_index[file_path]
            cache_file = self.cache_dir / "chunks" / f"{metadata.file_hash}.pkl"
            
            with open(cache_file, 'rb') as f:
                chunks_data = pickle.load(f)
                return [DocumentChunk(**chunk) for chunk in chunks_data]
                
        except Exception as e:
            self.logger.error(f"加载缓存块失败: {e}")
            return []
    
    def _save_chunks_to_cache(self, file_path: str, chunks: List[DocumentChunk], 
                             metadata: DocumentMetadata):
        """保存文档块到缓存"""
        if not self.enable_cache:
            return
        
        try:
            # 保存块数据
            cache_file = self.cache_dir / "chunks" / f"{metadata.file_hash}.pkl"
            chunks_data = [asdict(chunk) for chunk in chunks]
            
            with open(cache_file, 'wb') as f:
                pickle.dump(chunks_data, f)
            
            # 更新索引
            self.document_index[file_path] = metadata
            self._save_document_index()
            
            print(f"✅ 文档缓存已保存: {len(chunks)} 个块")
            
        except Exception as e:
            self.logger.error(f"保存缓存失败: {e}")
    
    def _clean_cache(self):
        """清理缓存"""
        if not self.enable_cache:
            return
        
        try:
            # 计算缓存大小
            total_size = 0
            cache_files = []
            
            for cache_file in (self.cache_dir / "chunks").glob("*.pkl"):
                size = cache_file.stat().st_size
                total_size += size
                cache_files.append((cache_file, size, cache_file.stat().st_mtime))
            
            # 如果超过限制，删除最旧的文件
            if total_size > self.max_cache_size_mb * 1024 * 1024:
                cache_files.sort(key=lambda x: x[2])  # 按修改时间排序
                
                for cache_file, size, _ in cache_files:
                    if total_size <= self.max_cache_size_mb * 1024 * 1024:
                        break
                    
                    cache_file.unlink()
                    total_size -= size
                    print(f"🗑️ 清理缓存文件: {cache_file.name}")
                    
        except Exception as e:
            self.logger.error(f"清理缓存失败: {e}")
    
    def _process_document_optimized(self, file_path: str, extract_images: bool = False, extract_tables: bool = True, force_reprocess: bool = False) -> List[DocumentChunk]:
        """优化的文档处理方法"""
        # 检查缓存
        if not force_reprocess and self._is_document_cached(file_path):
            return self._load_cached_chunks(file_path)
        
        # 获取文件信息
        file_stat = os.stat(file_path)
        file_ext = Path(file_path).suffix.lower()
        
        # 检查格式支持
        if file_ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {file_ext}")
        
        start_time = time.time()
        
        # 根据文件类型处理
        if file_ext == '.pdf':
            # 优先使用多模态处理器处理PDF
            if self.multimodal_processor:
                chunks = self._process_pdf_with_multimodal(file_path)
            else:
                chunks = self._process_pdf_optimized(file_path, extract_images, extract_tables)
        elif file_ext in ['.docx']:
            chunks = self._process_docx(file_path)
        elif file_ext in ['.doc']:
            chunks = self._process_doc(file_path)
        elif file_ext in ['.txt']:
            chunks = self._process_txt(file_path)
        elif file_ext in ['.md', '.markdown']:
            chunks = self._process_markdown(file_path)
        elif file_ext in ['.csv']:
            chunks = self._process_csv(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            chunks = self._process_excel(file_path)
        elif file_ext in ['.pptx']:
            chunks = self._process_pptx(file_path)
        elif file_ext in ['.json']:
            chunks = self._process_json(file_path)
        elif file_ext in ['.xml']:
            chunks = self._process_xml(file_path)
        elif file_ext in ['.html', '.htm']:
            chunks = self._process_html(file_path)
        elif file_ext in ['.rtf']:
            chunks = self._process_rtf(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {file_ext}")
        
        processing_time = time.time() - start_time
        
        # 创建文档元数据
        file_hash = self._get_file_hash(file_path)
        metadata = DocumentMetadata(
            file_path=file_path,
            file_hash=file_hash,
            file_size=file_stat.st_size,
            last_modified=file_stat.st_mtime,
            processing_time=processing_time,
            chunk_count=len(chunks),
            content_type=self._detect_content_type(chunks),
            format_type=file_ext,
            created_at=datetime.now().isoformat()
        )
        
        # 保存到缓存
        self._save_chunks_to_cache(file_path, chunks, metadata)
        
        return chunks
    
    def process_document(self, file_path: str, force_reprocess: bool = False) -> List[DocumentChunk]:
        """处理文档"""
        print(f"📄 处理文档: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 检查是否使用缓存
        if not force_reprocess and self._is_document_cached(file_path):
            print("⚡ 使用缓存数据")
            return self._load_cached_chunks(file_path)
        
        # 获取文件信息
        file_stat = os.stat(file_path)
        file_ext = Path(file_path).suffix.lower()
        
        # 检查格式支持
        if file_ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {file_ext}")
        
        # 处理文档
        start_time = time.time()
        processor = self.supported_formats[file_ext]
        chunks = processor(file_path)
        processing_time = time.time() - start_time
        
        # 创建元数据
        metadata = DocumentMetadata(
            file_path=file_path,
            file_hash=self._get_file_hash(file_path),
            file_size=file_stat.st_size,
            last_modified=file_stat.st_mtime,
            processing_time=processing_time,
            chunk_count=len(chunks),
            content_type=self._detect_content_type(chunks),
            format_type=file_ext,
            created_at=datetime.now().isoformat()
        )
        
        # 保存到缓存
        self._save_chunks_to_cache(file_path, chunks, metadata)
        
        # 清理缓存
        self._clean_cache()
        
        print(f"✅ 文档处理完成: {len(chunks)} 个块, 耗时 {processing_time:.2f}s")
        return chunks
    
    def _process_pdf_with_multimodal(self, file_path: str) -> List[DocumentChunk]:
        """使用多模态处理器处理PDF"""
        try:
            print(f"🔍 使用多模态处理器处理PDF: {file_path}")
            result = self.multimodal_processor.process_pdf_with_multimodal(file_path, is_url=False)
            
            chunks = []
            texts_to_embed = result.get('texts_to_embed', [])
            metadata_to_embed = result.get('metadata_to_embed', [])
            
            for i, (text, meta) in enumerate(zip(texts_to_embed, metadata_to_embed)):
                if text.strip():
                    # 转换元数据格式
                    chunk_metadata = {
                        "page": meta.get('page', 1),
                        "source": file_path,
                        "type": meta.get('type', 'text'),
                        "multimodal_id": meta.get('multimodal_id', f"chunk_{i}")
                    }
                    
                    chunk_type = meta.get('type', 'text')
                    chunks.append(self._create_chunk(text, chunk_metadata, chunk_type))
            
            print(f"✅ 多模态PDF处理完成: {len(chunks)} 个块")
            return chunks
            
        except Exception as e:
            print(f"⚠️ 多模态PDF处理失败，回退到标准处理: {e}")
            # 回退到标准PDF处理
            return self._process_pdf_optimized(file_path, extract_images=True, extract_tables=True)
    
    def _detect_content_type(self, chunks: List[DocumentChunk]) -> str:
        """检测内容类型"""
        type_counts = {}
        for chunk in chunks:
            chunk_type = chunk.chunk_type
            type_counts[chunk_type] = type_counts.get(chunk_type, 0) + 1
        
        if not type_counts:
            return "unknown"
        
        return max(type_counts, key=type_counts.get)
    
    def _create_chunk(self, content: str, metadata: Dict[str, Any], 
                     chunk_type: str = "text") -> DocumentChunk:
        """创建文档块"""
        chunk_id = hashlib.md5(content.encode()).hexdigest()[:16]
        return DocumentChunk(
            chunk_id=chunk_id,
            content=content,
            metadata=metadata,
            chunk_type=chunk_type
        )
    
    def _process_pdf_optimized(self, file_path: str, extract_images: bool = False, extract_tables: bool = True) -> List[DocumentChunk]:
        """优化的PDF处理方法"""
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF未安装，无法处理PDF文件")
        
        chunks = []
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取文本（始终执行）
                text = page.get_text()
                if text.strip():
                    metadata = {
                        "page": page_num + 1,
                        "source": file_path,
                        "type": "text"
                    }
                    chunks.append(self._create_chunk(text, metadata, "text"))
                
                # 可选：提取图像
                if extract_images:
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n - pix.alpha < 4:  # 确保不是CMYK
                                metadata = {
                                    "page": page_num + 1,
                                    "source": file_path,
                                    "type": "image",
                                    "image_index": img_index
                                }
                                chunks.append(self._create_chunk(f"图像数据 (页面 {page_num + 1})", metadata, "image"))
                            pix = None
                        except Exception as e:
                            self.logger.warning(f"提取图像失败: {e}")
                
                # 可选：提取表格
                if extract_tables:
                    try:
                        tables = page.find_tables()
                        for table_index, table in enumerate(tables):
                            try:
                                table_data = table.extract()
                                table_text = "\n".join(["\t".join(row) for row in table_data])
                                metadata = {
                                    "page": page_num + 1,
                                    "source": file_path,
                                    "type": "table",
                                    "table_index": table_index
                                }
                                chunks.append(self._create_chunk(table_text, metadata, "table"))
                            except Exception as e:
                                self.logger.warning(f"提取表格失败: {e}")
                    except Exception:
                        # 如果表格提取失败，继续处理其他内容
                        pass
            
            doc.close()
            
        except Exception as e:
            raise Exception(f"PDF处理失败: {e}")
        
        return chunks
    
    # 各种格式的处理器
    def _process_pdf(self, file_path: str) -> List[DocumentChunk]:
        """处理PDF文件"""
        if not PYMUPDF_AVAILABLE:
            raise ImportError("PyMuPDF未安装，无法处理PDF文件")
        
        chunks = []
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取文本
                text = page.get_text()
                if text.strip():
                    metadata = {
                        "page": page_num + 1,
                        "source": file_path,
                        "type": "text"
                    }
                    chunks.append(self._create_chunk(text, metadata, "text"))
                
                # 提取图像
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n - pix.alpha < 4:  # 确保不是CMYK
                            img_data = pix.tobytes("png")
                            metadata = {
                                "page": page_num + 1,
                                "source": file_path,
                                "type": "image",
                                "image_index": img_index
                            }
                            chunks.append(self._create_chunk(f"图像数据 (页面 {page_num + 1})", metadata, "image"))
                        pix = None
                    except Exception as e:
                        self.logger.warning(f"提取图像失败: {e}")
                
                # 提取表格
                tables = page.find_tables()
                for table_index, table in enumerate(tables):
                    try:
                        table_data = table.extract()
                        table_text = "\n".join(["\t".join(row) for row in table_data])
                        metadata = {
                            "page": page_num + 1,
                            "source": file_path,
                            "type": "table",
                            "table_index": table_index
                        }
                        chunks.append(self._create_chunk(table_text, metadata, "table"))
                    except Exception as e:
                        self.logger.warning(f"提取表格失败: {e}")
            
            doc.close()
            
        except Exception as e:
            raise Exception(f"PDF处理失败: {e}")
        
        return chunks
    
    def _process_docx(self, file_path: str) -> List[DocumentChunk]:
        """处理DOCX文件"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx未安装，无法处理DOCX文件")
        
        chunks = []
        try:
            doc = Document(file_path)
            
            # 处理段落
            for para_index, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    metadata = {
                        "paragraph": para_index + 1,
                        "source": file_path,
                        "type": "text"
                    }
                    chunks.append(self._create_chunk(paragraph.text, metadata, "text"))
            
            # 处理表格
            for table_index, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                table_text = "\n".join(["\t".join(row) for row in table_data])
                metadata = {
                    "table": table_index + 1,
                    "source": file_path,
                    "type": "table"
                }
                chunks.append(self._create_chunk(table_text, metadata, "table"))
                
        except Exception as e:
            raise Exception(f"DOCX处理失败: {e}")
        
        return chunks
    
    def _process_doc(self, file_path: str) -> List[DocumentChunk]:
        """处理DOC文件"""
        # 对于DOC文件，可以使用python-docx2txt或其他库
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            metadata = {
                "source": file_path,
                "type": "text"
            }
            return [self._create_chunk(text, metadata, "text")]
        except ImportError:
            raise ImportError("docx2txt未安装，无法处理DOC文件")
    
    def _process_txt(self, file_path: str) -> List[DocumentChunk]:
        """处理TXT文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 按段落分割
            paragraphs = content.split('\n\n')
            for i, paragraph in enumerate(paragraphs):
                if paragraph.strip():
                    metadata = {
                        "paragraph": i + 1,
                        "source": file_path,
                        "type": "text"
                    }
                    chunks.append(self._create_chunk(paragraph, metadata, "text"))
                    
        except Exception as e:
            raise Exception(f"TXT处理失败: {e}")
        
        return chunks
    
    def _process_markdown(self, file_path: str) -> List[DocumentChunk]:
        """处理Markdown文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 简单的Markdown解析
            lines = content.split('\n')
            current_section = []
            section_level = 0
            
            for line in lines:
                if line.startswith('#'):
                    # 保存当前段落
                    if current_section:
                        section_text = '\n'.join(current_section)
                        metadata = {
                            "section_level": section_level,
                            "source": file_path,
                            "type": "text"
                        }
                        chunks.append(self._create_chunk(section_text, metadata, "text"))
                        current_section = []
                    
                    # 开始新段落
                    section_level = len(line) - len(line.lstrip('#'))
                    current_section = [line]
                else:
                    current_section.append(line)
            
            # 保存最后一个段落
            if current_section:
                section_text = '\n'.join(current_section)
                metadata = {
                    "section_level": section_level,
                    "source": file_path,
                    "type": "text"
                }
                chunks.append(self._create_chunk(section_text, metadata, "text"))
                
        except Exception as e:
            raise Exception(f"Markdown处理失败: {e}")
        
        return chunks
    
    def _process_csv(self, file_path: str) -> List[DocumentChunk]:
        """处理CSV文件"""
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas未安装，无法处理CSV文件")
        
        chunks = []
        try:
            df = pd.read_csv(file_path)
            
            # 将整个CSV作为一个表格块
            csv_text = df.to_string(index=False)
            metadata = {
                "rows": len(df),
                "columns": len(df.columns),
                "source": file_path,
                "type": "table"
            }
            chunks.append(self._create_chunk(csv_text, metadata, "table"))
            
            # 如果数据量大，可以按行分割
            if len(df) > 100:
                for i in range(0, len(df), 50):
                    chunk_df = df.iloc[i:i+50]
                    chunk_text = chunk_df.to_string(index=False)
                    metadata = {
                        "rows": len(chunk_df),
                        "columns": len(chunk_df.columns),
                        "chunk_start": i,
                        "source": file_path,
                        "type": "table"
                    }
                    chunks.append(self._create_chunk(chunk_text, metadata, "table"))
                    
        except Exception as e:
            raise Exception(f"CSV处理失败: {e}")
        
        return chunks
    
    def _process_excel(self, file_path: str) -> List[DocumentChunk]:
        """处理Excel文件"""
        if not PANDAS_AVAILABLE or not EXCEL_AVAILABLE:
            raise ImportError("pandas或openpyxl未安装，无法处理Excel文件")
        
        chunks = []
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    sheet_text = df.to_string(index=False)
                    metadata = {
                        "sheet": sheet_name,
                        "rows": len(df),
                        "columns": len(df.columns),
                        "source": file_path,
                        "type": "table"
                    }
                    chunks.append(self._create_chunk(sheet_text, metadata, "table"))
                    
        except Exception as e:
            raise Exception(f"Excel处理失败: {e}")
        
        return chunks
    
    def _process_pptx(self, file_path: str) -> List[DocumentChunk]:
        """处理PPTX文件"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx未安装，无法处理PPTX文件")
        
        chunks = []
        try:
            prs = Presentation(file_path)
            
            for slide_index, slide in enumerate(prs.slides):
                slide_text = []
                
                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content = '\n'.join(slide_text)
                    metadata = {
                        "slide": slide_index + 1,
                        "source": file_path,
                        "type": "text"
                    }
                    chunks.append(self._create_chunk(text_content, metadata, "text"))
                    
        except Exception as e:
            raise Exception(f"PPTX处理失败: {e}")
        
        return chunks
    
    def _process_ppt(self, file_path: str) -> List[DocumentChunk]:
        """处理PPT文件"""
        # PPT文件处理较复杂，可能需要使用win32com或其他库
        raise NotImplementedError("PPT格式处理尚未实现")
    
    def _process_json(self, file_path: str) -> List[DocumentChunk]:
        """处理JSON文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # 将JSON转换为可读文本
            json_text = json.dumps(data, ensure_ascii=False, indent=2)
            metadata = {
                "source": file_path,
                "type": "data"
            }
            chunks.append(self._create_chunk(json_text, metadata, "text"))
            
        except Exception as e:
            raise Exception(f"JSON处理失败: {e}")
        
        return chunks
    
    def _process_xml(self, file_path: str) -> List[DocumentChunk]:
        """处理XML文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            metadata = {
                "source": file_path,
                "type": "data"
            }
            chunks.append(self._create_chunk(content, metadata, "text"))
            
        except Exception as e:
            raise Exception(f"XML处理失败: {e}")
        
        return chunks
    
    def _process_html(self, file_path: str) -> List[DocumentChunk]:
        """处理HTML文件"""
        chunks = []
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            text = soup.get_text()
            
            metadata = {
                "source": file_path,
                "type": "text"
            }
            chunks.append(self._create_chunk(text, metadata, "text"))
            
        except ImportError:
            raise ImportError("BeautifulSoup4未安装，无法处理HTML文件")
        except Exception as e:
            raise Exception(f"HTML处理失败: {e}")
        
        return chunks
    
    def _process_rtf(self, file_path: str) -> List[DocumentChunk]:
        """处理RTF文件"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            metadata = {
                "source": file_path,
                "type": "text"
            }
            return [self._create_chunk(text, metadata, "text")]
            
        except ImportError:
            raise ImportError("striprtf未安装，无法处理RTF文件")
        except Exception as e:
            raise Exception(f"RTF处理失败: {e}")
    
    def batch_process_documents(self, file_paths: List[str], 
                               max_workers: int = 4,
                               show_progress: bool = True,
                               extract_images: bool = False,
                               extract_tables: bool = True,
                               force_reprocess: bool = False) -> Dict[str, List[DocumentChunk]]:
        """批量处理文档 - 优化版本"""
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import time
        
        start_time = time.time()
        results = {}
        total_files = len(file_paths)
        completed_files = 0
        total_chunks = 0
        
        print(f"🚀 开始批量处理 {total_files} 个文档...")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交任务
            future_to_path = {
                executor.submit(self._process_document_optimized, path, extract_images, extract_tables, force_reprocess): path 
                for path in file_paths
            }
            
            # 收集结果
            for future in as_completed(future_to_path):
                file_path = future_to_path[future]
                try:
                    chunks = future.result()
                    results[file_path] = chunks
                    completed_files += 1
                    total_chunks += len(chunks)
                    
                    if show_progress:
                        progress = (completed_files / total_files) * 100
                        print(f"📄 [{completed_files}/{total_files}] ({progress:.1f}%) 完成: {os.path.basename(file_path)} - {len(chunks)} 个块")
                        
                except Exception as e:
                    print(f"❌ 处理失败: {os.path.basename(file_path)} - {e}")
                    results[file_path] = []
                    completed_files += 1
        
        processing_time = time.time() - start_time
        print(f"\n✅ 批量处理完成!")
        print(f"📊 统计信息:")
        print(f"   - 处理文档: {total_files} 个")
        print(f"   - 生成块数: {total_chunks} 个")
        print(f"   - 总耗时: {processing_time:.2f}s")
        print(f"   - 平均每文档: {processing_time/total_files:.2f}s")
        print(f"   - 平均每块: {processing_time/max(total_chunks, 1):.3f}s")
        
        return results
    
    def get_document_info(self, file_path: str) -> Optional[DocumentMetadata]:
        """获取文档信息"""
        return self.document_index.get(file_path)
    
    def list_cached_documents(self) -> List[DocumentMetadata]:
        """列出所有缓存的文档"""
        return list(self.document_index.values())
    
    def clear_cache(self, file_path: Optional[str] = None):
        """清理缓存"""
        if file_path:
            # 清理特定文档的缓存
            if file_path in self.document_index:
                metadata = self.document_index[file_path]
                cache_file = self.cache_dir / "chunks" / f"{metadata.file_hash}.pkl"
                if cache_file.exists():
                    cache_file.unlink()
                del self.document_index[file_path]
                self._save_document_index()
                print(f"🗑️ 已清理缓存: {file_path}")
            else:
                print(f"⚠️ 文档未在缓存中: {file_path}")
        else:
            # 清理所有缓存
            import shutil
            if self.cache_dir.exists():
                shutil.rmtree(self.cache_dir)
                self.cache_dir.mkdir(exist_ok=True)
                (self.cache_dir / "metadata").mkdir(exist_ok=True)
                (self.cache_dir / "chunks").mkdir(exist_ok=True)
                (self.cache_dir / "embeddings").mkdir(exist_ok=True)
            self.document_index = {}
            print("🗑️ 已清理所有缓存")
    
    def clear_all_cache_interactive(self) -> bool:
        """交互式清理所有缓存"""
        stats = self.get_cache_stats()
        if stats['cached_documents'] == 0:
            print("ℹ️ 当前没有缓存的文档")
            return False
        
        print(f"\n📊 当前缓存状态:")
        print(f"   - 缓存文档数: {stats['cached_documents']} 个")
        print(f"   - 缓存文件数: {stats['cache_files']} 个")
        print(f"   - 缓存大小: {stats['total_size_mb']:.2f} MB")
        print(f"   - 使用率: {stats['cache_usage_percent']:.1f}%")
        
        try:
            choice = input("\n🗑️ 确认清理所有缓存? (y/N): ").strip().lower()
            if choice == 'y':
                self.clear_cache()
                print("✅ 缓存清理完成")
                return True
            else:
                print("❌ 取消清理操作")
                return False
        except KeyboardInterrupt:
            print("\n❌ 操作被用户取消")
            return False
    
    def remove_document_from_cache(self, file_path: str) -> bool:
        """从缓存中移除特定文档"""
        if file_path in self.document_index:
            self.clear_cache(file_path)
            return True
        else:
            print(f"⚠️ 文档未在缓存中: {file_path}")
            return False
    
    def get_cache_stats(self) -> Dict[str, Any]:
        """获取缓存统计信息"""
        total_size = 0
        file_count = 0
        
        if self.cache_dir.exists():
            for cache_file in (self.cache_dir / "chunks").glob("*.pkl"):
                total_size += cache_file.stat().st_size
                file_count += 1
        
        return {
            "cached_documents": len(self.document_index),
            "cache_files": file_count,
            "total_size_mb": total_size / (1024 * 1024),
            "max_size_mb": self.max_cache_size_mb,
            "cache_usage_percent": (total_size / (1024 * 1024)) / self.max_cache_size_mb * 100
        }

def main():
    """测试函数"""
    manager = EnhancedDocumentManager()
    
    # 测试处理文档
    test_file = "test.pdf"  # 替换为实际文件路径
    if os.path.exists(test_file):
        chunks = manager.process_document(test_file)
        print(f"处理结果: {len(chunks)} 个块")
        
        # 显示缓存统计
        stats = manager.get_cache_stats()
        print(f"缓存统计: {stats}")

if __name__ == "__main__":
    main()